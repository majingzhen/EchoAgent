from __future__ import annotations

import io
import os
from datetime import datetime
from typing import Any


def _fmt_date(dt: datetime | str | None) -> str:
    if dt is None:
        return ""
    if isinstance(dt, str):
        return dt[:19].replace("T", " ")
    return dt.strftime("%Y-%m-%d %H:%M")


def _resolve_font() -> str:
    """Try to find a CJK font path, return empty string if not found."""
    candidates = [
        "C:/Windows/Fonts/simhei.ttf",
        "C:/Windows/Fonts/msyh.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/System/Library/Fonts/PingFang.ttc",
    ]
    for p in candidates:
        if os.path.exists(p):
            return p
    return ""


_FONT_PATH = _resolve_font()


class ExportService:
    def export_pdf(self, module: str, title: str, data: dict[str, Any]) -> bytes:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

        buf = io.BytesIO()
        doc = SimpleDocTemplate(
            buf, pagesize=A4,
            leftMargin=2 * cm, rightMargin=2 * cm,
            topMargin=2 * cm, bottomMargin=2 * cm,
        )

        font_name = "Helvetica"
        if _FONT_PATH:
            try:
                pdfmetrics.registerFont(TTFont("CJK", _FONT_PATH))
                font_name = "CJK"
            except Exception:
                pass

        title_style = ParagraphStyle("T", fontName=font_name, fontSize=18, spaceAfter=8)
        meta_style = ParagraphStyle("M", fontName=font_name, fontSize=11, spaceAfter=12, textColor="#829ab1")
        h2_style = ParagraphStyle("H2", fontName=font_name, fontSize=14, spaceBefore=14, spaceAfter=6)
        body_style = ParagraphStyle("B", fontName=font_name, fontSize=11, leading=18)

        story: list[Any] = []
        story.append(Paragraph(title, title_style))
        story.append(Paragraph(f"生成时间：{_fmt_date(datetime.utcnow())}", meta_style))
        story.append(Spacer(1, 0.3 * cm))

        for section_title, items in self._get_sections(module, data):
            story.append(Paragraph(section_title, h2_style))
            for item in items:
                story.append(Paragraph(f"• {str(item)}", body_style))
            story.append(Spacer(1, 0.2 * cm))

        doc.build(story)
        return buf.getvalue()

    def export_pptx(self, module: str, title: str, data: dict[str, Any]) -> bytes:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        body_ph = slide.placeholders[1]
        body_ph.text = f"生成时间：{_fmt_date(datetime.utcnow())}"

        # Content slides (one per section)
        for section_title, items in self._get_sections(module, data):
            if not items:
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = section_title
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, item in enumerate(items[:10]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(item)
                p.level = 0
                p.runs[0].font.size = Pt(16)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def _get_sections(self, module: str, data: dict[str, Any]) -> list[tuple[str, list[str]]]:
        handlers = {
            "simulation": self._simulation_sections,
            "workshop": self._workshop_sections,
            "sentiment-guard": self._sentiment_guard_sections,
            "strategy-advisor": self._strategy_advisor_sections,
            "focus-group": self._focus_group_sections,
        }
        fn = handlers.get(module)
        return fn(data) if fn else []

    def _simulation_sections(self, data: dict[str, Any]) -> list[tuple[str, list[str]]]:
        sections = []
        if data.get("executive_summary"):
            sections.append(("执行摘要", [data["executive_summary"]]))
        metrics = data.get("metrics") or {}
        if metrics:
            items = [f"{k}: {v}" for k, v in metrics.items() if not isinstance(v, (dict, list))]
            if items:
                sections.append(("关键指标", items))
        if data.get("risks"):
            sections.append(("风险点", [str(r) for r in data["risks"]]))
        if data.get("suggestions"):
            sections.append(("优化建议", [str(s) for s in data["suggestions"]]))
        if data.get("sample_comments"):
            sections.append(("典型评论", [str(c) for c in data["sample_comments"][:5]]))
        if not sections:
            # Fallback to status info
            items = []
            if data.get("platform"):
                items.append(f"平台：{data['platform']}")
            if data.get("status"):
                items.append(f"状态：{data['status']}")
            if data.get("total_rounds"):
                items.append(f"总轮次：{data['total_rounds']}")
            if items:
                sections.append(("模拟概况", items))
        return sections

    def _workshop_sections(self, data: dict[str, Any]) -> list[tuple[str, list[str]]]:
        sections = []
        brief = data.get("brief") or {}
        if isinstance(brief, dict):
            text = brief.get("text") or brief.get("description") or str(brief)
        else:
            text = str(brief)
        if text and text not in ("{}", ""):
            sections.append(("创作简报", [text[:500]]))
        payload = data.get("payload") or {}
        if payload.get("final_content"):
            sections.append(("最终内容", [str(payload["final_content"])[:600]]))
        elif payload.get("versions"):
            items = [f"版本 {i+1}：{str(v)[:200]}" for i, v in enumerate(payload["versions"][:3])]
            if items:
                sections.append(("内容版本", items))
        if data.get("insights"):
            sections.append(("市场洞察", [str(i) for i in data["insights"][:8]]))
        return sections

    def _sentiment_guard_sections(self, data: dict[str, Any]) -> list[tuple[str, list[str]]]:
        sections = []
        if data.get("event_description"):
            sections.append(("事件描述", [data["event_description"]]))
        payload = data.get("payload") or {}
        risk = payload.get("risk_assessment") or {}
        if risk:
            items = [f"{k}: {v}" for k, v in risk.items() if isinstance(v, str)][:6]
            if items:
                sections.append(("风险评估", items))
        if payload.get("best_plan"):
            sections.append(("最优应对方案", [str(payload["best_plan"])]))
        if payload.get("summary"):
            sections.append(("总结", [str(payload["summary"])]))
        plans = payload.get("response_plans") or []
        if plans:
            items = []
            for i, p in enumerate(plans[:5]):
                if isinstance(p, dict):
                    name = p.get("name", f"方案{i+1}")
                    desc = p.get("description", "")[:120]
                    items.append(f"{name}：{desc}")
            if items:
                sections.append(("应对方案", items))
        return sections

    def _strategy_advisor_sections(self, data: dict[str, Any]) -> list[tuple[str, list[str]]]:
        sections = []
        if data.get("question"):
            sections.append(("咨询问题", [data["question"]]))
        if data.get("context_info"):
            sections.append(("背景信息", [data["context_info"][:300]]))
        payload = data.get("payload") or {}
        for key, label in [("advice", "策略建议"), ("analysis", "分析"), ("action_plan", "行动计划")]:
            val = payload.get(key)
            if not val:
                continue
            if isinstance(val, list):
                sections.append((label, [str(v) for v in val[:8]]))
            else:
                sections.append((label, [str(val)[:600]]))
        return sections

    def _focus_group_sections(self, data: dict[str, Any]) -> list[tuple[str, list[str]]]:
        sections = []
        if data.get("topic"):
            sections.append(("讨论话题", [data["topic"]]))
        summary = data.get("summary") or {}
        for key, label in [("consensus", "共识观点"), ("divergence", "分歧点"), ("key_insights", "关键洞察"), ("recommendations", "建议")]:
            val = summary.get(key)
            if val:
                sections.append((label, [str(v) for v in val]))
        messages = data.get("messages") or []
        persona_msgs = [m for m in messages if m.get("sender_type") == "persona"]
        if persona_msgs:
            items = [f"{m.get('persona_name', '画像')}：{m.get('content', '')[:150]}" for m in persona_msgs[:6]]
            sections.append(("典型发言", items))
        return sections
